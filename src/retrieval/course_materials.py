from __future__ import annotations

import json
import math
import re
from collections import Counter
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Any


SUPPORTED_EXTENSIONS = {".docx", ".md", ".pdf", ".pptx", ".txt"}
TOKEN_PATTERN = re.compile(r"[\w\u0600-\u06FF]+", re.UNICODE)
MIN_EXTRACTED_TEXT_LENGTH = 20
SEARCH_STOPWORDS = {
    "a",
    "about",
    "an",
    "and",
    "are",
    "describe",
    "explain",
    "for",
    "from",
    "in",
    "is",
    "lecture",
    "material",
    "materials",
    "me",
    "my",
    "notes",
    "of",
    "on",
    "please",
    "summarise",
    "summarize",
    "tell",
    "the",
    "this",
    "to",
    "what",
}
ACRONYM_ALIASES = {
    "grc": [
        "governance risk compliance",
        "governance, risk, and compliance",
        "governance risk and compliance",
    ],
}


@dataclass(frozen=True)
class RetrievedChunk:
    source_name: str
    section: str
    text: str
    score: float
    chunk_index: int = 1
    course_id: str | None = None
    course_name: str | None = None
    file_name: str | None = None
    page_or_chunk: str | None = None

    def citation(self) -> str:
        if self.course_name:
            return f"{self.course_name} - {self.source_name} - {self.section}/chunk {self.chunk_index}"
        return f"{self.source_name} ({self.section})"


class CourseMaterialIndexer:
    """Extracts uploaded course files into a small persistent vector store."""

    def __init__(
        self,
        *,
        uploads_dir: Path,
        vector_store_dir: Path,
        chunk_size: int = 900,
        chunk_overlap: int = 140,
    ) -> None:
        self.uploads_dir = uploads_dir
        self.vector_store_dir = vector_store_dir
        self.store_path = vector_store_dir / "course_materials.json"
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap

    def index_all(self, *, course_id: str | None = None, course_name: str | None = None) -> dict[str, Any]:
        self.uploads_dir.mkdir(parents=True, exist_ok=True)
        results = []
        search_root = self._course_upload_dir(course_id) if course_id else self.uploads_dir
        for file_path in sorted(search_root.glob("*")):
            if file_path.is_file() and file_path.suffix.lower() in SUPPORTED_EXTENSIONS:
                results.append(self.index_file(file_path, course_id=course_id, course_name=course_name))

        stats = self.stats(course_id=course_id)
        return {
            "ok": all(result["ok"] for result in results) if results else True,
            "indexed_files": sum(1 for result in results if result["ok"]),
            "errors": [result for result in results if not result["ok"]],
            "stats": stats,
        }

    def index_file(
        self,
        file_path: Path,
        *,
        course_id: str | None = None,
        course_name: str | None = None,
    ) -> dict[str, Any]:
        file_path = file_path.resolve()
        if not file_path.exists() or not file_path.is_file():
            return {"ok": False, "file_name": file_path.name, "reason": "File does not exist."}
        if file_path.suffix.lower() not in SUPPORTED_EXTENSIONS:
            return {"ok": False, "file_name": file_path.name, "reason": "Unsupported file type."}

        store = self._load_store()
        existing_chunks = [
            chunk
            for chunk in store["chunks"]
            if chunk.get("source_path") == str(file_path) and chunk.get("course_id") == course_id
        ]
        stat = file_path.stat()
        if existing_chunks and all(
            chunk.get("source_mtime") == stat.st_mtime and chunk.get("source_size") == stat.st_size
            for chunk in existing_chunks
        ):
            return {"ok": True, "file_name": file_path.name, "chunks": len(existing_chunks), "skipped": True}

        try:
            sections = self._extract_sections(file_path)
        except Exception as exc:  # pragma: no cover - defensive path for optional parsers
            return {"ok": False, "file_name": file_path.name, "reason": str(exc)}

        extracted_text_length = sum(len(str(section.get("text") or "").strip()) for section in sections)
        text_previews = [
            self._preview_text(str(section.get("text") or ""))
            for section in sections[:3]
            if str(section.get("text") or "").strip()
        ]
        if extracted_text_length < MIN_EXTRACTED_TEXT_LENGTH:
            return {
                "ok": False,
                "file_name": file_path.name,
                "reason": "Little or no extractable text was found.",
                "extraction_status": "too_little_text",
                "extracted_text_length": extracted_text_length,
                "text_previews": text_previews,
                "chunks": 0,
            }

        chunks = []
        for section in sections:
            for chunk_index, chunk_text in enumerate(self._chunk_text(section["text"]), start=1):
                vector = self._embed(chunk_text)
                if not vector:
                    continue
                page_or_chunk = section["label"] or f"chunk {chunk_index}"
                chunks.append(
                    {
                        "id": f"{course_id or 'legacy'}:{file_path.name}:{section['label']}:{chunk_index}",
                        "course_id": course_id,
                        "course_name": course_name,
                        "file_name": file_path.name,
                        "source_name": file_path.name,
                        "source_path": str(file_path),
                        "source_mtime": stat.st_mtime,
                        "source_size": stat.st_size,
                        "page_or_chunk": page_or_chunk,
                        "section": section["label"],
                        "chunk_index": chunk_index,
                        "text": chunk_text,
                        "embedding": vector,
                        "indexed_at_utc": datetime.utcnow().isoformat(timespec="seconds"),
                    }
                )

        if not chunks:
            return {
                "ok": False,
                "file_name": file_path.name,
                "reason": "Little or no extractable text was found.",
                "extraction_status": "too_little_text",
                "extracted_text_length": extracted_text_length,
                "text_previews": text_previews,
                "chunks": 0,
            }

        store["chunks"] = [
            chunk
            for chunk in store["chunks"]
            if not (chunk.get("source_path") == str(file_path) and chunk.get("course_id") == course_id)
        ]
        store["chunks"].extend(chunks)
        store["updated_at_utc"] = datetime.utcnow().isoformat(timespec="seconds")
        self._save_store(store)

        return {
            "ok": True,
            "file_name": file_path.name,
            "chunks": len(chunks),
            "extracted_text_length": extracted_text_length,
            "text_previews": text_previews,
        }

    def remove_file(
        self,
        file_path: Path,
        *,
        delete_source: bool = True,
        course_id: str | None = None,
    ) -> dict[str, Any]:
        file_path = file_path.resolve()
        uploads_dir = self.uploads_dir.resolve()
        if not file_path.is_relative_to(uploads_dir):
            return {"ok": False, "file_name": file_path.name, "reason": "File is outside the uploads directory."}

        store = self._load_store()
        previous_count = len(store["chunks"])
        store["chunks"] = [
            chunk
            for chunk in store["chunks"]
            if not (
                chunk.get("source_path") == str(file_path)
                and (course_id is None or chunk.get("course_id") == course_id)
            )
        ]
        removed_chunks = previous_count - len(store["chunks"])

        file_deleted = False
        if delete_source and file_path.exists():
            if not file_path.is_file():
                return {"ok": False, "file_name": file_path.name, "reason": "Path is not a file."}
            file_path.unlink()
            file_deleted = True

        store["updated_at_utc"] = datetime.utcnow().isoformat(timespec="seconds")
        self._save_store(store)

        return {
            "ok": True,
            "file_name": file_path.name,
            "file_deleted": file_deleted,
            "removed_chunks": removed_chunks,
        }

    def rename_course(self, *, course_id: str, course_name: str) -> dict[str, Any]:
        store = self._load_store()
        renamed_chunks = 0
        for chunk in store["chunks"]:
            if chunk.get("course_id") == course_id:
                chunk["course_name"] = course_name
                renamed_chunks += 1

        if renamed_chunks:
            store["updated_at_utc"] = datetime.utcnow().isoformat(timespec="seconds")
            self._save_store(store)

        return {"ok": True, "renamed_chunks": renamed_chunks}

    def remove_course(self, *, course_id: str) -> dict[str, Any]:
        store = self._load_store()
        previous_count = len(store["chunks"])
        store["chunks"] = [chunk for chunk in store["chunks"] if chunk.get("course_id") != course_id]
        removed_chunks = previous_count - len(store["chunks"])

        course_dir = self._course_upload_dir(course_id)
        removed_files = 0
        if course_dir.exists():
            for file_path in course_dir.glob("*"):
                if file_path.is_file():
                    file_path.unlink()
                    removed_files += 1
            try:
                course_dir.rmdir()
            except OSError:
                pass

        if removed_chunks or removed_files:
            store["updated_at_utc"] = datetime.utcnow().isoformat(timespec="seconds")
            self._save_store(store)

        return {
            "ok": True,
            "removed_chunks": removed_chunks,
            "removed_files": removed_files,
        }

    def search(
        self,
        query: str,
        *,
        top_k: int = 4,
        min_score: float = 0.05,
        course_id: str | None = None,
    ) -> list[RetrievedChunk]:
        query_vector = self._embed(query)
        query_terms = self._query_terms(query)
        query_phrases = self._query_phrases(query)
        if not query_vector and not query_terms and not query_phrases:
            return []

        store = self._load_store()
        scored_chunks = []
        for chunk in store["chunks"]:
            if course_id is not None and chunk.get("course_id") != course_id:
                continue
            score = max(
                self._cosine_similarity(query_vector, chunk.get("embedding", {})),
                self._lexical_match_score(chunk, query_terms=query_terms, query_phrases=query_phrases),
            )
            if score >= min_score:
                scored_chunks.append(
                    RetrievedChunk(
                        source_name=chunk.get("source_name") or chunk.get("file_name") or "",
                        section=chunk.get("section") or chunk.get("page_or_chunk") or "",
                        text=chunk.get("text") or "",
                        score=round(score, 4),
                        chunk_index=chunk.get("chunk_index", 1),
                        course_id=chunk.get("course_id"),
                        course_name=chunk.get("course_name"),
                        file_name=chunk.get("file_name") or chunk.get("source_name"),
                        page_or_chunk=chunk.get("page_or_chunk") or chunk.get("section"),
                    )
                )

        return sorted(scored_chunks, key=lambda item: item.score, reverse=True)[:top_k]

    def topic_suggestions(self, *, course_id: str | None = None, limit: int = 3) -> list[str]:
        """Return lightweight topic hints from indexed chunks in one course scope."""
        store = self._load_store()
        candidates: list[str] = []
        for chunk in store["chunks"]:
            if course_id is not None and chunk.get("course_id") != course_id:
                continue
            source_stem = Path(str(chunk.get("source_name") or "")).stem
            if source_stem:
                candidates.append(source_stem)
            candidates.extend(self._keyword_suggestions(str(chunk.get("text") or "")))

        suggestions = []
        seen: set[str] = set()
        for candidate in candidates:
            normalized = self._clean_topic(candidate)
            key = normalized.casefold()
            if not normalized or key in seen:
                continue
            suggestions.append(normalized)
            seen.add(key)
            if len(suggestions) >= limit:
                break
        return suggestions

    def chunks_for_source(
        self,
        source_name: str,
        *,
        course_id: str | None = None,
        top_k: int | None = None,
    ) -> list[RetrievedChunk]:
        source_name = str(source_name or "").strip()
        if not source_name:
            return []

        store = self._load_store()
        chunks = []
        for chunk in store["chunks"]:
            if course_id is not None and chunk.get("course_id") != course_id:
                continue
            if chunk.get("source_name") != source_name:
                continue
            chunks.append(
                RetrievedChunk(
                    source_name=chunk["source_name"],
                    section=chunk["section"],
                    text=chunk["text"],
                    score=1.0,
                    chunk_index=chunk.get("chunk_index", 1),
                    course_id=chunk.get("course_id"),
                    course_name=chunk.get("course_name"),
                    file_name=chunk.get("file_name") or chunk.get("source_name"),
                    page_or_chunk=chunk.get("page_or_chunk") or chunk.get("section"),
                )
            )

        ordered = sorted(chunks, key=lambda item: (item.section, item.chunk_index))
        if top_k is None:
            return ordered
        return ordered[: max(int(top_k), 0)]

    def stats(self, *, course_id: str | None = None) -> dict[str, Any]:
        store = self._load_store()
        chunks = [
            chunk
            for chunk in store["chunks"]
            if course_id is None or chunk.get("course_id") == course_id
        ]
        source_names = sorted({chunk["source_name"] for chunk in chunks})
        return {
            "files": len(source_names),
            "chunks": len(chunks),
            "sources": source_names,
            "updated_at_utc": store.get("updated_at_utc"),
        }

    def diagnose_course(
        self,
        *,
        course_id: str | None,
        course_name: str | None = None,
        query: str | None = None,
        preview_limit: int = 3,
    ) -> dict[str, Any]:
        """Return internal diagnostics for active-course upload and retrieval state."""
        self.uploads_dir.mkdir(parents=True, exist_ok=True)
        upload_root = self._course_upload_dir(course_id) if course_id else self.uploads_dir
        uploaded_files = [
            file_path
            for file_path in sorted(upload_root.glob("*"))
            if file_path.is_file() and file_path.suffix.lower() in SUPPORTED_EXTENSIONS
        ]

        extraction = []
        for file_path in uploaded_files:
            try:
                sections = self._extract_sections(file_path)
            except Exception as exc:  # pragma: no cover - diagnostics should not break chat
                extraction.append(
                    {
                        "file_name": file_path.name,
                        "extracted_text_length": 0,
                        "text_previews": [],
                        "error": str(exc),
                    }
                )
                continue
            previews = [
                self._preview_text(str(section.get("text") or ""))
                for section in sections[:preview_limit]
                if str(section.get("text") or "").strip()
            ]
            extraction.append(
                {
                    "file_name": file_path.name,
                    "extracted_text_length": sum(len(str(section.get("text") or "").strip()) for section in sections),
                    "text_previews": previews,
                }
            )

        store = self._load_store()
        chunks = [
            chunk
            for chunk in store["chunks"]
            if course_id is None or chunk.get("course_id") == course_id
        ]
        query_terms = self._query_terms(query or "")
        query_phrases = self._query_phrases(query or "")
        chunks_containing_query = [
            {
                "file_name": chunk.get("file_name") or chunk.get("source_name"),
                "page_or_chunk": chunk.get("page_or_chunk") or chunk.get("section"),
                "preview": self._preview_text(str(chunk.get("text") or "")),
            }
            for chunk in chunks
            if self._lexical_match_score(chunk, query_terms=query_terms, query_phrases=query_phrases) > 0
        ]

        return {
            "active_course_id": course_id,
            "active_course_name": course_name,
            "uploaded_file_count": len(uploaded_files),
            "uploaded_file_names": [file_path.name for file_path in uploaded_files],
            "extraction": extraction,
            "indexed_chunk_count": len(chunks),
            "indexed_sources": sorted({chunk.get("file_name") or chunk.get("source_name") for chunk in chunks}),
            "query": query,
            "query_terms": sorted(query_terms),
            "query_phrases": sorted(query_phrases),
            "chunks_containing_query_count": len(chunks_containing_query),
            "chunks_containing_query": chunks_containing_query[:preview_limit],
        }

    def _extract_sections(self, file_path: Path) -> list[dict[str, str]]:
        suffix = file_path.suffix.lower()
        if suffix in {".txt", ".md"}:
            return [{"label": "text", "text": self._read_text_file(file_path)}]
        if suffix == ".pdf":
            return self._extract_pdf(file_path)
        if suffix == ".docx":
            return self._extract_docx(file_path)
        if suffix == ".pptx":
            return self._extract_pptx(file_path)
        return []

    def _extract_pdf(self, file_path: Path) -> list[dict[str, str]]:
        try:
            import fitz
        except ImportError as exc:  # pragma: no cover - depends on local environment
            raise RuntimeError("Install PyMuPDF to index PDF files.") from exc

        sections = []
        with fitz.open(file_path) as document:
            for page_index, page in enumerate(document, start=1):
                text = page.get_text("text").strip()
                if text:
                    sections.append({"label": f"page {page_index}", "text": text})
        return sections

    def _extract_docx(self, file_path: Path) -> list[dict[str, str]]:
        try:
            from docx import Document
        except ImportError as exc:  # pragma: no cover - depends on local environment
            raise RuntimeError("Install python-docx to index DOCX files.") from exc

        document = Document(file_path)
        text = "\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text.strip())
        return [{"label": "document", "text": text}]

    def _extract_pptx(self, file_path: Path) -> list[dict[str, str]]:
        try:
            from pptx import Presentation
        except ImportError as exc:  # pragma: no cover - depends on local environment
            raise RuntimeError("Install python-pptx to index PPTX files.") from exc

        presentation = Presentation(file_path)
        sections = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())
            if text_parts:
                sections.append({"label": f"slide {slide_index}", "text": "\n".join(text_parts)})
        return sections

    def _read_text_file(self, file_path: Path) -> str:
        for encoding in ("utf-8", "utf-8-sig", "cp1252"):
            try:
                return file_path.read_text(encoding=encoding)
            except UnicodeDecodeError:
                continue
        return file_path.read_text(encoding="utf-8", errors="ignore")

    def _preview_text(self, text: str, *, max_length: int = 180) -> str:
        compact = re.sub(r"\s+", " ", str(text or "")).strip()
        if len(compact) <= max_length:
            return compact
        return compact[:max_length].rsplit(" ", 1)[0].rstrip() + "..."

    def _chunk_text(self, text: str) -> list[str]:
        normalized = re.sub(r"\s+", " ", text).strip()
        if not normalized:
            return []
        if len(normalized) <= self.chunk_size:
            return [normalized]

        chunks = []
        start = 0
        while start < len(normalized):
            end = min(start + self.chunk_size, len(normalized))
            if end < len(normalized):
                sentence_break = normalized.rfind(". ", start, end)
                if sentence_break > start + int(self.chunk_size * 0.55):
                    end = sentence_break + 1
            chunk = normalized[start:end].strip()
            if chunk:
                chunks.append(chunk)
            if end >= len(normalized):
                break
            start = max(end - self.chunk_overlap, start + 1)
        return chunks

    def _embed(self, text: str) -> dict[str, float]:
        tokens = [token.lower() for token in TOKEN_PATTERN.findall(text)]
        if not tokens:
            return {}
        counts = Counter(tokens)
        norm = math.sqrt(sum(value * value for value in counts.values()))
        return {token: round(value / norm, 6) for token, value in counts.items()}

    def _cosine_similarity(self, left: dict[str, float], right: dict[str, float]) -> float:
        if not left or not right:
            return 0.0
        smaller, larger = (left, right) if len(left) <= len(right) else (right, left)
        return sum(weight * larger.get(token, 0.0) for token, weight in smaller.items())

    def _normalize_search_text(self, text: str) -> str:
        tokens = [token.casefold() for token in TOKEN_PATTERN.findall(str(text or ""))]
        return " ".join(tokens)

    def _query_terms(self, query: str) -> set[str]:
        tokens = [token.casefold() for token in TOKEN_PATTERN.findall(str(query or ""))]
        terms = {token for token in tokens if token not in SEARCH_STOPWORDS}
        if not terms:
            terms = set(tokens)
        for token in list(terms):
            for alias in ACRONYM_ALIASES.get(token, []):
                terms.update(TOKEN_PATTERN.findall(alias.casefold()))
        return {term for term in terms if term and term not in SEARCH_STOPWORDS}

    def _query_phrases(self, query: str) -> set[str]:
        normalized = self._normalize_search_text(query)
        phrases = {normalized} if normalized else set()
        for token in TOKEN_PATTERN.findall(str(query or "").casefold()):
            phrases.add(token)
            for alias in ACRONYM_ALIASES.get(token, []):
                phrases.add(self._normalize_search_text(alias))
        return {phrase for phrase in phrases if phrase}

    def _lexical_match_score(
        self,
        chunk: dict[str, Any],
        *,
        query_terms: set[str],
        query_phrases: set[str],
    ) -> float:
        if not query_terms and not query_phrases:
            return 0.0
        searchable = " ".join(
            str(chunk.get(field) or "")
            for field in ("text", "source_name", "file_name", "section", "page_or_chunk")
        )
        normalized = self._normalize_search_text(searchable)
        if not normalized:
            return 0.0

        chunk_terms = set(normalized.split())
        for phrase in query_phrases:
            if not phrase:
                continue
            phrase_terms = phrase.split()
            if len(phrase_terms) == 1 and phrase in chunk_terms:
                return 1.0
            if len(phrase_terms) > 1 and phrase in normalized:
                return 1.0

        exact_matches = query_terms & chunk_terms
        if exact_matches:
            return max(0.35, min(0.95, len(exact_matches) / max(len(query_terms), 1)))

        partial_matches = [
            term
            for term in query_terms
            if len(term) >= 4 and any(term in chunk_term or chunk_term in term for chunk_term in chunk_terms)
        ]
        if partial_matches:
            return max(0.12, min(0.5, len(partial_matches) / max(len(query_terms), 1)))
        return 0.0

    def _load_store(self) -> dict[str, Any]:
        if not self.store_path.exists():
            return {"version": 1, "updated_at_utc": None, "chunks": []}
        try:
            return json.loads(self.store_path.read_text(encoding="utf-8"))
        except (json.JSONDecodeError, OSError):
            return {"version": 1, "updated_at_utc": None, "chunks": []}

    def _save_store(self, store: dict[str, Any]) -> None:
        self.vector_store_dir.mkdir(parents=True, exist_ok=True)
        self.store_path.write_text(json.dumps(store, indent=2, ensure_ascii=False), encoding="utf-8")

    def _course_upload_dir(self, course_id: str | None) -> Path:
        if not course_id:
            return self.uploads_dir
        return self.uploads_dir / course_id

    def _keyword_suggestions(self, text: str) -> list[str]:
        tokens = [token for token in TOKEN_PATTERN.findall(text) if len(token) > 2]
        stopwords = {
            "and",
            "are",
            "for",
            "from",
            "into",
            "that",
            "the",
            "this",
            "through",
            "using",
            "with",
        }
        suggestions = []
        for index in range(len(tokens)):
            window = tokens[index : index + 3]
            if len(window) < 2:
                continue
            if all(token.casefold() in stopwords for token in window):
                continue
            phrase = " ".join(window)
            if any(token.casefold() not in stopwords for token in window):
                suggestions.append(phrase)
        return suggestions

    def _clean_topic(self, value: str) -> str:
        cleaned = re.sub(r"^\d{8}_\d{6}_", "", str(value or ""))
        cleaned = re.sub(r"[_\-]+", " ", cleaned)
        cleaned = re.sub(r"\s+", " ", cleaned).strip(" .,:;")
        return cleaned[:80]
